"""Generate Ekspor PPT from SIG template (Black/Light) + export dashboard data.

Reads a JSON payload from stdin: {records, months, theme}. `theme` is "black"
(dark slides) or "light" (light slides) and selects both the SIG template and the
colour palette. Output path is argv[1].
"""
import sys, os, json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn

HERE = os.path.dirname(os.path.abspath(__file__))
TEMPLATES = {
    "black": os.path.join(HERE, "ppt-templates", "sig-black.pptx"),
    "light": os.path.join(HERE, "ppt-templates", "sig-light.pptx"),
}
FONT = "SIG Text"

# SIG accent palette (theme1.xml) — shared across both themes for chart series.
SERIES = [
    RGBColor(0x15, 0x60, 0x82), RGBColor(0xE9, 0x71, 0x32),
    RGBColor(0x19, 0x6B, 0x24), RGBColor(0x0F, 0x9E, 0xD5),
    RGBColor(0xA0, 0x2B, 0x93), RGBColor(0x4E, 0xA7, 0x2E),
    RGBColor(0x38, 0xBD, 0xF8), RGBColor(0xFB, 0xBF, 0x24),
]
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Content & Section slides are light-background in BOTH SIG templates, so body
# text is always dark. Only the Title slide bg (and KPI-card styling) differ by
# theme; the Ending slide is rendered by the template itself.
CONTENT = {
    "text":   RGBColor(0x0F, 0x17, 0x2A),
    "muted":  RGBColor(0x47, 0x55, 0x69),
    "accent": RGBColor(0x15, 0x60, 0x82),
    "grid":   RGBColor(0xE2, 0xE8, 0xF0),
    "header": RGBColor(0x15, 0x60, 0x82),
    "row":    RGBColor(0xF1, 0xF5, 0xF9),
    "axis":   RGBColor(0x33, 0x41, 0x55),
}
C = CONTENT

THEMES = {
    # Black template: dark title slide → light title text; bold dark-navy cards.
    "black": {
        "title_text": RGBColor(0xF5, 0xF7, 0xFA),
        "subtitle":   RGBColor(0xCB, 0xD5, 0xE1),
        "box":        RGBColor(0x0E, 0x28, 0x41),
        "box_line":   RGBColor(0x24, 0x44, 0x66),
        "box_label":  RGBColor(0x94, 0xA3, 0xB8),
        "box_value":  RGBColor(0x38, 0xBD, 0xF8),
    },
    # Light template: light title slide → dark title text; soft grey cards.
    "light": {
        "title_text": RGBColor(0x0F, 0x17, 0x2A),
        "subtitle":   RGBColor(0x47, 0x55, 0x69),
        "box":        RGBColor(0xF1, 0xF5, 0xF9),
        "box_line":   RGBColor(0xCB, 0xD5, 0xE1),
        "box_label":  RGBColor(0x64, 0x74, 0x8B),
        "box_value":  RGBColor(0x15, 0x60, 0x82),
    },
}

PAL = THEMES["black"]  # set in generate()


# ----------------------------------------------------------------------------
# low-level helpers
# ----------------------------------------------------------------------------
def sf(run, size=10, bold=False, color=None):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def atxt(slide, text, x, y, w, h, size=10, bold=False, color=None, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    sf(r, size, bold, color if color is not None else C["text"])
    return tb


def get_ph(slide, idx):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def set_ph_text(slide, idx, text, size, bold, color, align=None):
    ph = get_ph(slide, idx)
    if ph is None:
        return None
    ph.text = ""
    p = ph.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = text
    sf(r, size, bold, color)
    if align is not None:
        p.alignment = align
    return ph


def delete_all_slides(prs):
    """Drop the template's built-in sample slides so output starts clean."""
    lst = prs.slides._sldIdLst
    for sid in list(lst):
        prs.part.drop_rel(sid.get(qn("r:id")))
        lst.remove(sid)


def _no_fill():
    return parse_xml('<c:spPr %s><a:noFill/><a:ln><a:noFill/></a:ln></c:spPr>' % nsdecls("c", "a"))


def style_chart(chart, legend=False, value_axis=True):
    """Make a native chart blend into the (image) slide background and use
    dark, readable text on the light content slides."""
    cs = chart._chartSpace
    chart_el = cs.find(qn("c:chart"))
    # transparent chart area + plot area so the slide background shows through
    chart_el.addnext(_no_fill())
    plot_area = chart_el.find(qn("c:plotArea"))
    if plot_area is not None:
        plot_area.append(_no_fill())

    chart.has_title = False
    chart.has_legend = legend
    if legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(9)
        chart.legend.font.color.rgb = C["text"]
        chart.legend.font.name = FONT

    try:
        cat = chart.category_axis
        cat.tick_labels.font.size = Pt(8)
        cat.tick_labels.font.color.rgb = C["axis"]
        cat.tick_labels.font.name = FONT
        cat.format.line.color.rgb = C["grid"]
        cat.has_major_gridlines = False
    except Exception:
        pass
    try:
        val = chart.value_axis
        val.visible = value_axis
        val.has_major_gridlines = False  # keep charts clean / gridline-free
        if value_axis:
            val.tick_labels.font.size = Pt(7)
            val.tick_labels.font.color.rgb = C["axis"]
            val.tick_labels.font.name = FONT
            val.format.line.color.rgb = C["grid"]
    except Exception:
        pass


def add_bar(slide, labels, values, x, y, w, h, color=None):
    """Horizontal bars, no value axis / data labels — the adjacent table carries
    the exact figures, so bars stay clean."""
    cd = CategoryChartData()
    cd.categories = labels
    cd.add_series("Value", values)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, w, h, cd).chart
    style_chart(chart, legend=False, value_axis=False)
    plot = chart.plots[0]
    plot.gap_width = 60
    plot.has_data_labels = False
    plot.series[0].format.fill.solid()
    plot.series[0].format.fill.fore_color.rgb = color or SERIES[0]


def add_pie(slide, labels, values, x, y, w, h):
    cd = CategoryChartData()
    cd.categories = labels
    cd.add_series("Value", values)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, w, h, cd).chart
    style_chart(chart, legend=True, value_axis=False)
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.show_percentage = True
    plot.data_labels.show_value = False
    plot.data_labels.number_format = "0%"
    plot.data_labels.number_format_is_linked = False
    plot.data_labels.font.size = Pt(8)
    plot.data_labels.font.color.rgb = WHITE
    plot.data_labels.font.name = FONT
    for i, pt in enumerate(plot.series[0].points):
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = SERIES[i % len(SERIES)]


def add_line(slide, labels, values, x, y, w, h, color=None):
    cd = CategoryChartData()
    cd.categories = labels
    cd.add_series("Value", values)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, x, y, w, h, cd).chart
    style_chart(chart, legend=False, value_axis=False)
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.number_format = '#,##0,,"M"'
    plot.data_labels.number_format_is_linked = False
    plot.data_labels.font.size = Pt(8)
    plot.data_labels.font.color.rgb = C["axis"]
    plot.data_labels.font.name = FONT
    s = plot.series[0]
    s.format.line.color.rgb = color or SERIES[0]
    s.format.line.width = Pt(2.5)
    s.smooth = True


def add_line_dual(slide, labels, vals1, vals2, name1, name2, x, y, w, h):
    cd = CategoryChartData()
    cd.categories = labels
    cd.add_series(name1, vals1)
    cd.add_series(name2, vals2)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, x, y, w, h, cd).chart
    style_chart(chart, legend=True, value_axis=False)
    colors = [SERIES[0], SERIES[2]]
    for plot in chart.plots:
        plot.has_data_labels = False
        for si, s in enumerate(plot.series):
            s.format.line.color.rgb = colors[si % len(colors)]
            s.format.line.width = Pt(2.5)
            s.smooth = True


def add_table(slide, headers, rows, x, y, w, rh=Inches(0.26)):
    nr = len(rows) + 1
    nc = len(headers)
    ts = slide.shapes.add_table(nr, nc, x, y, w, rh * nr)
    tbl = ts.table
    for ci, htext in enumerate(headers):
        c = tbl.cell(0, ci)
        c.text = ""
        r = c.text_frame.paragraphs[0].add_run()
        r.text = htext
        sf(r, 9, True, WHITE)
        c.fill.solid()
        c.fill.fore_color.rgb = C["header"]
    for ri, row in enumerate(rows):
        for ci, v in enumerate(row):
            c = tbl.cell(ri + 1, ci)
            c.text = ""
            r = c.text_frame.paragraphs[0].add_run()
            r.text = str(v)
            sf(r, 8, False, C["text"])
            c.fill.solid()
            c.fill.fore_color.rgb = C["row"]
    return ts


def add_kpi_box(slide, label, value, sub, x, y, w, h, vc=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = PAL["box"]
    shape.line.color.rgb = PAL["box_line"]
    shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    atxt(slide, label, x + Inches(0.15), y + Inches(0.10), w - Inches(0.3), Inches(0.3), 9, False, PAL["box_label"])
    atxt(slide, value, x + Inches(0.15), y + Inches(0.40), w - Inches(0.3), Inches(0.4), 16, True, vc or PAL["box_value"])
    if sub:
        atxt(slide, sub, x + Inches(0.15), y + Inches(0.92), w - Inches(0.3), Inches(0.25), 8, False, PAL["box_label"])


# ----------------------------------------------------------------------------
# slide scaffolding
# ----------------------------------------------------------------------------
def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_ph_text(slide, 0, title, 30, True, PAL["title_text"])
    # Place the subtitle as our own textbox just under the title for reliable
    # position and contrast (template body placeholder sits in a busy corner).
    atxt(slide, subtitle, Inches(0.55), Inches(3.1), Inches(7.5), Inches(0.5),
         13, False, PAL["subtitle"])
    return slide


def add_section_slide(prs, title):
    # The section background is a busy architectural image; lay a full-width SIG
    # blue band behind the title so white text stays high-contrast.
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    ph = get_ph(slide, 0)
    if ph is not None:
        ph.text = ""
    band = slide.shapes.add_shape(1, Inches(0), Inches(3.05), Inches(13.333), Inches(1.4))
    band.fill.solid()
    band.fill.fore_color.rgb = SERIES[0]
    band.line.fill.background()
    band.shadow.inherit = False
    tf = band.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    sf(r, 26, True, WHITE)
    return slide


def content_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[2])  # Title Only
    set_ph_text(slide, 14, title, 18, True, C["accent"])
    return slide


def add_narrative_slide(prs, narrative):
    """Optional AI-generated 'Ringkasan Eksekutif' slide. No-op if narrative is empty."""
    if not narrative:
        return
    summary = (narrative.get("executiveSummary") or "").strip()
    insights = [str(s).strip() for s in (narrative.get("insights") or []) if str(s).strip()]
    if not summary and not insights:
        return

    slide = content_slide(prs, "RINGKASAN EKSEKUTIF")
    box = slide.shapes.add_textbox(Inches(0.55), Inches(1.3), Inches(12.2), Inches(5.6))
    tf = box.text_frame
    tf.word_wrap = True

    first = True
    if summary:
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = summary
        sf(r, 14, False, C["text"])
        p.space_after = Pt(16)
        first = False

    for ins in insights[:5]:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run()
        r.text = f"•  {ins}"
        sf(r, 12, False, C["muted"])
        p.space_after = Pt(10)
    return slide


def add_end_slide(prs):
    # Ending layout already carries the branded "Terima Kasih" design.
    prs.slides.add_slide(prs.slide_layouts[7])


# ----------------------------------------------------------------------------
# formatting
# ----------------------------------------------------------------------------
def fusd(v):
    a = abs(v)
    if a >= 1e9:
        return f"USD {a/1e9:.2f} B"
    if a >= 1e6:
        return f"USD {a/1e6:.2f} M"
    return f"USD {a:,.0f}"


def fidr(v):
    a = abs(v)
    if a >= 1e12:
        return f"Rp {a/1e12:.2f} T"
    if a >= 1e9:
        return f"Rp {a/1e9:.2f} M"
    return f"Rp {a:,.0f}"


def fton(v):
    a = abs(v)
    if a >= 1e6:
        return f"{a/1e6:.2f}M MT"
    if a >= 1e3:
        return f"{a/1e3:.0f}K MT"
    return f"{v:,.0f} MT"


def pct(v):
    return f"{v*100:.0f}%"


def sv(items, key="value"):
    return sorted(items, key=lambda x: x.get(key, 0), reverse=True)


def is_paid(r):
    return bool(r.get("paymentDate")) or str(r.get("paymentStatus", "")).lower() == "paid"


# ----------------------------------------------------------------------------
# main
# ----------------------------------------------------------------------------
def generate(input_json):
    global PAL
    data = json.loads(input_json)
    records = data["records"]
    months = data.get("months", [])
    theme = (data.get("theme") or "black").lower()
    if theme not in TEMPLATES:
        theme = "black"
    PAL = THEMES[theme]

    latest = months[0]["label"] if months else "N/A"
    rate_est = 16500

    ts = sum(r["usdValue"] for r in records)
    paid = [r for r in records if is_paid(r)]
    unpaid = [r for r in records if not is_paid(r)]
    tp = sum(r["usdValue"] for r in paid)
    t_out = sum(r["usdValue"] for r in unpaid)
    t_ton = sum(r["tonnage"] for r in records)

    # Company breakdown
    cm = {}
    for r in records:
        cc = r["companyCode"] or "Unmapped"
        e = cm.get(cc, {"sales": 0, "paid": 0, "outstanding": 0, "tonnage": 0, "count": 0})
        e["sales"] += r["usdValue"]; e["tonnage"] += r["tonnage"]; e["count"] += 1
        if is_paid(r): e["paid"] += r["usdValue"]
        else: e["outstanding"] += r["usdValue"]
        cm[cc] = e
    companies = sv([{"code": k, **v} for k, v in cm.items()], "sales")

    # Destination
    dm = {}
    for r in records:
        d = r.get("destination", "") or "Unmapped"
        e = dm.get(d, {"value": 0, "tonnage": 0, "count": 0})
        e["value"] += r["usdValue"]; e["tonnage"] += r["tonnage"]; e["count"] += 1
        dm[d] = e
    destinations = sv([{"label": k, **v} for k, v in dm.items()])

    # Buyer
    bm = {}
    for r in records:
        b = r.get("buyer", "") or "Unmapped"
        e = bm.get(b, {"value": 0, "count": 0})
        e["value"] += r["usdValue"]; e["count"] += 1
        bm[b] = e
    buyers = sv([{"label": k, **v} for k, v in bm.items()])

    # Monthly
    mm = {}
    for r in records:
        pk = r.get("periodKey", "")
        e = mm.get(pk, {"label": r.get("periodLabel", ""), "sales": 0, "paid": 0, "tonnage": 0})
        e["sales"] += r["usdValue"]; e["tonnage"] += r["tonnage"]
        if is_paid(r): e["paid"] += r["usdValue"]
        mm[pk] = e
    monthly = sorted([{"key": k, **v} for k, v in mm.items()], key=lambda x: x["key"])

    # Aging (unpaid)
    am = {}
    for r in unpaid:
        b = r.get("agingBucket", "") or "Unmapped"
        e = am.get(b, {"value": 0, "count": 0})
        e["value"] += r["usdValue"]; e["count"] += 1
        am[b] = e
    aging = sv([{"label": k, **v} for k, v in am.items()])

    # Demurrage / FX difference
    dem = [r for r in records if r.get("paymentRate", 0) or r.get("exchangeDifference", 0)]
    total_dem_impact = sum(abs(r.get("exchangeImpact", 0) or 0) for r in dem)

    prs = Presentation(TEMPLATES[theme])
    delete_all_slides(prs)

    # 1. Title
    add_title_slide(
        prs,
        f"LAPORAN KINERJA EKSPOR\n{latest}",
        f"UNIT of COMMERCIAL FINANCE 2  |  {len(records):,} transaksi",
    )

    # 1b. AI executive summary (only if narrative was generated)
    add_narrative_slide(prs, data.get("narrative"))

    # 2. OVERVIEW
    slide = content_slide(prs, "OVERVIEW — Ringkasan Kinerja Ekspor")
    kpis = [
        ("Total Penjualan (USD)", fusd(ts), f"{len(records):,} transaksi"),
        ("Total Penerimaan (USD)", fusd(tp), f"{len(paid):,} paid"),
        ("Piutang Outstanding (USD)", fusd(t_out), f"{len(unpaid):,} unpaid"),
        ("Total Tonase", fton(t_ton), ""),
        ("Estimasi IDR (Penjualan)", fidr(ts * rate_est), f"Rate ~Rp {rate_est:,}"),
        ("Estimasi IDR (Outstanding)", fidr(t_out * rate_est), ""),
    ]
    for i, (lb, vl, sb) in enumerate(kpis):
        col, row = i % 3, i // 3
        add_kpi_box(slide, lb, vl, sb, Inches(0.5 + col * 4.2), Inches(1.3 + row * 1.65), Inches(3.9), Inches(1.4))

    # 3. RKAP (no budget data → realisasi bulanan)
    slide = content_slide(prs, "RKAP — Realisasi Penjualan Bulanan")
    ml = [m["label"] for m in monthly]
    msales = [m["sales"] for m in monthly]
    if len(monthly) > 1:
        atxt(slide, "Realisasi Penjualan Bulanan (USD)", Inches(0.5), Inches(1.2), Inches(6), Inches(0.3), 10, False, C["muted"])
        add_line(slide, ml, msales, Inches(0.5), Inches(1.6), Inches(12.3), Inches(5.0))
    else:
        atxt(slide, f"Data tersedia untuk 1 periode ({latest}): {fusd(ts)}", Inches(1), Inches(3), Inches(10), Inches(0.5), 16, False, C["text"])

    # 4. TREN EKSPOR
    slide = content_slide(prs, "TREN EKSPOR — Penjualan vs Penerimaan")
    if len(monthly) > 1:
        atxt(slide, "Penjualan vs Penerimaan Bulanan (USD)", Inches(0.5), Inches(1.2), Inches(6), Inches(0.3), 10, False, C["muted"])
        add_line_dual(slide, ml, msales, [m["paid"] for m in monthly], "Penjualan", "Penerimaan",
                      Inches(0.5), Inches(1.6), Inches(12.3), Inches(5.0))
    else:
        atxt(slide, "Data multi-periode belum tersedia.", Inches(1), Inches(3), Inches(10), Inches(0.5), 16, False, C["text"])

    # 5. Section break
    add_section_slide(prs, "EKSPOR SIG GROUP")

    # 6. EKSPOR SIG GROUP — per entity (one slide: bar + table)
    slide = content_slide(prs, "EKSPOR SIG GROUP — Per Entitas")
    if companies:
        atxt(slide, "Penjualan per Entitas (USD)", Inches(0.5), Inches(1.2), Inches(6), Inches(0.3), 10, False, C["muted"])
        add_bar(slide, [c["code"] for c in companies[:8]], [c["sales"] for c in companies[:8]],
                Inches(0.5), Inches(1.6), Inches(6.4), Inches(5.0))
        hdrs = ["Entitas", "Penjualan", "Penerimaan", "Outstanding", "Tonase"]
        rows = [[c["code"], fusd(c["sales"]), fusd(c["paid"]), fusd(c["outstanding"]), fton(c["tonnage"])]
                for c in companies[:10]]
        add_table(slide, hdrs, rows, Inches(7.2), Inches(1.6), Inches(5.6))

    # 7. TUJUAN EKSPOR
    slide = content_slide(prs, "TUJUAN EKSPOR — Negara Tujuan")
    if destinations:
        td = destinations[:10]
        atxt(slide, "Top 10 Negara Tujuan (USD)", Inches(0.5), Inches(1.2), Inches(6), Inches(0.3), 10, False, C["muted"])
        add_bar(slide, [d["label"] for d in td], [d["value"] for d in td], Inches(0.5), Inches(1.6), Inches(6.6), Inches(5.0))
        tdv = sum(d["value"] for d in destinations) or 1
        hdrs = ["Negara", "Nilai (USD)", "Tonase", "Share"]
        rows = [[d["label"], fusd(d["value"]), fton(d.get("tonnage", 0)), pct(d["value"] / tdv)] for d in td]
        add_table(slide, hdrs, rows, Inches(7.4), Inches(1.6), Inches(5.4))

    # 8. PROGNOSA
    slide = content_slide(prs, "PROGNOSA — Proyeksi Penerimaan")
    ppm = {}
    for r in unpaid:
        pp = (r.get("plannedPaymentDate") or "").strip()
        # Only bucket by month when the value is an ISO-ish YYYY-MM date; messy
        # free-text planned dates fall back to TBD instead of bogus labels.
        head = pp[:7]
        pmonth = head if len(head) == 7 and head[:4].isdigit() and head[4] == "-" and head[5:7].isdigit() else "TBD"
        e = ppm.get(pmonth, {"value": 0, "count": 0})
        e["value"] += r["usdValue"]; e["count"] += 1
        ppm[pmonth] = e
    prognosa = sorted([{"label": k, **v} for k, v in ppm.items()], key=lambda x: x["label"])
    if prognosa:
        atxt(slide, "Proyeksi Penerimaan per Rencana Bayar (USD)", Inches(0.5), Inches(1.2), Inches(7), Inches(0.3), 10, False, C["muted"])
        add_bar(slide, [p["label"] for p in prognosa], [p["value"] for p in prognosa], Inches(0.5), Inches(1.6), Inches(6.6), Inches(5.0))
        tpv = sum(p["value"] for p in prognosa) or 1
        hdrs = ["Periode", "Nilai (USD)", "Transaksi", "Share"]
        rows = [[p["label"], fusd(p["value"]), str(p["count"]), pct(p["value"] / tpv)] for p in prognosa]
        add_table(slide, hdrs, rows, Inches(7.4), Inches(1.6), Inches(5.4))
    else:
        atxt(slide, "Belum ada data rencana pembayaran.", Inches(1), Inches(3), Inches(10), Inches(0.5), 16, False, C["text"])

    # 9. DEMURRAGE-DESPATCH
    slide = content_slide(prs, "DEMURRAGE - DESPATCH SIG GROUP")
    if dem:
        add_kpi_box(slide, "Transaksi dgn Selisih Kurs", f"{len(dem):,}", "", Inches(0.5), Inches(1.3), Inches(3.9), Inches(1.4))
        add_kpi_box(slide, "Total Dampak Selisih (IDR)", fidr(total_dem_impact), "", Inches(4.6), Inches(1.3), Inches(3.9), Inches(1.4))
        dcm = {}
        for r in dem:
            cc = r["companyCode"] or "Unmapped"
            e = dcm.get(cc, {"value": 0, "count": 0})
            e["value"] += abs(r.get("exchangeImpact", 0) or 0); e["count"] += 1
            dcm[cc] = e
        dc = sv([{"label": k, **v} for k, v in dcm.items()])[:8]
        if dc:
            atxt(slide, "Dampak Selisih Kurs per Entitas (IDR)", Inches(0.5), Inches(3.0), Inches(6), Inches(0.3), 10, False, C["muted"])
            add_bar(slide, [d["label"] for d in dc], [d["value"] for d in dc], Inches(0.5), Inches(3.3), Inches(6.2), Inches(3.8))
            hdrs = ["Entitas", "Transaksi", "Dampak (IDR)"]
            rows = [[d["label"], str(d["count"]), fidr(d["value"])] for d in dc]
            add_table(slide, hdrs, rows, Inches(7.2), Inches(3.3), Inches(5.6))
    else:
        atxt(slide, "Tidak ada data selisih kurs / demurrage.", Inches(1), Inches(3), Inches(10), Inches(0.5), 16, False, C["text"])

    # 10. (extra) PIUTANG AGING
    if aging:
        slide = content_slide(prs, "PIUTANG AGING — Aging Bucket Outstanding")
        atxt(slide, "Distribusi Piutang per Aging (USD)", Inches(0.5), Inches(1.2), Inches(6), Inches(0.3), 10, False, C["muted"])
        add_pie(slide, [a["label"] for a in aging], [a["value"] for a in aging], Inches(0.5), Inches(1.6), Inches(5.6), Inches(4.8))
        tav = sum(a["value"] for a in aging) or 1
        hdrs = ["Aging", "Nilai (USD)", "Transaksi", "Share"]
        rows = [[a["label"], fusd(a["value"]), str(a["count"]), pct(a["value"] / tav)] for a in aging]
        add_table(slide, hdrs, rows, Inches(6.6), Inches(1.6), Inches(6.2))

    # 11. (extra) TOP BUYERS
    if buyers:
        slide = content_slide(prs, "TOP BUYERS EKSPOR")
        tb = buyers[:10]
        atxt(slide, "Top 10 Buyers (USD)", Inches(0.5), Inches(1.2), Inches(6), Inches(0.3), 10, False, C["muted"])
        add_bar(slide, [b["label"] for b in tb], [b["value"] for b in tb], Inches(0.5), Inches(1.6), Inches(6.6), Inches(5.0))
        tbv = sum(b["value"] for b in buyers) or 1
        hdrs = ["Buyer", "Nilai (USD)", "Transaksi", "Share"]
        rows = [[b["label"], fusd(b["value"]), str(b["count"]), pct(b["value"] / tbv)] for b in tb]
        add_table(slide, hdrs, rows, Inches(7.4), Inches(1.6), Inches(5.4))

    # 12. Ending
    add_end_slide(prs)
    return prs


if __name__ == "__main__":
    in_json = sys.stdin.read()
    prs = generate(in_json)
    out = sys.argv[1] if len(sys.argv) > 1 else "output_ekspor.pptx"
    prs.save(out)
    print(f"OK:{out}")
