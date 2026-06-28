"""Generate PPTX from SIG template + dashboard data."""
import sys, json, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn

# Use the repo-tracked template (same file as the export deck) so generation is
# portable and never depends on the temporary .openclaw-attachments folder.
TEMPLATE = os.path.join(os.path.dirname(os.path.abspath(__file__)), "ppt-templates", "sig-black.pptx")

# ── Helpers ──

def format_currency(value, compact=False):
    if compact:
        abs_v = abs(value)
        if abs_v >= 1_000_000_000_000:
            return f"Rp {abs_v/1_000_000_000_000:.1f} T"
        if abs_v >= 1_000_000_000:
            return f"Rp {abs_v/1_000_000_000:.1f} M"
        if abs_v >= 1_000_000:
            return f"Rp {abs_v/1_000_000:.0f} jt"
    return f"Rp {int(round(value)):,}".replace(",", ".")

def format_pct(value):
    return f"{value * 100:.0f}%"

def today_label():
    from datetime import datetime
    months = ["Januari","Februari","Maret","April","Mei","Juni",
              "Juli","Agustus","September","Oktober","November","Desember"]
    now = datetime.now()
    return f"{now.day:02d} {months[now.month-1]} {now.year}"

CHART_COLORS = [
    RGBColor(0x38, 0xBD, 0xF8),  # sky
    RGBColor(0x81, 0x8C, 0xF8),  # indigo
    RGBColor(0x34, 0xD3, 0x99),  # emerald
    RGBColor(0xFB, 0xBF, 0x24),  # amber
    RGBColor(0xEF, 0x44, 0x44),  # red
    RGBColor(0xFB, 0x92, 0x3C),  # orange
    RGBColor(0xA7, 0x8B, 0xFA),  # violet
    RGBColor(0xF4, 0x72, 0xB6),  # pink
]

FONT = "SIG Text"
FONT_FALLBACK = "Arial"

def set_font(run, size_pt=10, bold=False, color=None, name=FONT):
    run.font.name = name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color

# ── Slide builders ──

def delete_all_slides(prs):
    """Drop the template's built-in sample slides so output starts clean."""
    lst = prs.slides._sldIdLst
    for sid in list(lst):
        prs.part.drop_rel(sid.get(qn("r:id")))
        lst.remove(sid)

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0:
            shape.text = ""
            p = shape.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = title
            set_font(run, 28, True, RGBColor(0xFF, 0xFF, 0xFF))
        elif shape.placeholder_format.idx == 10:
            shape.text = ""
            p = shape.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = subtitle
            set_font(run, 14, False, RGBColor(0x94, 0xA3, 0xB8))
    return slide

def add_section_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[4])  # Section Header
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0:
            shape.text = ""
            p = shape.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = title
            set_font(run, 24, True, RGBColor(0xFF, 0xFF, 0xFF))
    return slide

def add_content_slide(prs, slide_title):
    """Add a slide using Title Only layout with a manually added title bar."""
    slide = prs.slides.add_slide(prs.slide_layouts[2])  # Title Only
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 14:
            shape.text = ""
            p = shape.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = slide_title
            set_font(run, 18, True, RGBColor(0x38, 0xBD, 0xF8))
    return slide

def add_narrative_slide(prs, narrative):
    """Optional AI-generated 'Ringkasan Eksekutif' slide. No-op if narrative is empty."""
    if not narrative:
        return
    summary = (narrative.get("executiveSummary") or "").strip()
    insights = [str(s).strip() for s in (narrative.get("insights") or []) if str(s).strip()]
    if not summary and not insights:
        return

    slide = add_content_slide(prs, "Ringkasan Eksekutif")
    box = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(12.1), Inches(5.4))
    tf = box.text_frame
    tf.word_wrap = True

    first = True
    if summary:
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = summary
        set_font(r, 13, False, RGBColor(0xE2, 0xE8, 0xF0))
        p.space_after = Pt(14)
        first = False

    for ins in insights[:5]:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run()
        r.text = f"•  {ins}"
        set_font(r, 12, False, RGBColor(0x94, 0xA3, 0xB8))
        p.space_after = Pt(8)
    return slide

def add_end_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[7])  # Ending Slide
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 3:
            shape.text = ""
            p = shape.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = "TERIMA KASIH"
            set_font(run, 28, True, RGBColor(0x38, 0xBD, 0xF8))
    return slide

# ── Add chart helpers ──

def add_pie_chart(slide, chart_title, labels, values, left, top, width, height):
    """Add a pie chart to a slide."""
    chart_data = CategoryChartData()
    chart_data.categories = labels
    chart_data.add_series("Value", values)
    
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, left, top, width, height, chart_data
    ).chart
    
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(8)
    chart.legend.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
    
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.show_percentage = True
    data_labels.show_value = False
    data_labels.font.size = Pt(8)
    data_labels.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    
    # Color the slices
    for i, point in enumerate(plot.series[0].points):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = CHART_COLORS[i % len(CHART_COLORS)]
    
    return chart

def add_bar_chart(slide, chart_title, labels, values, left, top, width, height):
    """Add a horizontal bar chart."""
    chart_data = CategoryChartData()
    chart_data.categories = labels
    chart_data.add_series("Amount", values)
    
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    ).chart
    
    chart.has_legend = False
    chart.category_axis.tick_labels.font.size = Pt(8)
    chart.category_axis.tick_labels.font.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    chart.value_axis.tick_labels.font.size = Pt(7)
    chart.value_axis.tick_labels.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
    
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.font.size = Pt(8)
    plot.data_labels.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = RGBColor(0x38, 0xBD, 0xF8)
    
    return chart

def add_line_chart(slide, chart_title, labels, values, left, top, width, height):
    """Add a line chart."""
    chart_data = CategoryChartData()
    chart_data.categories = labels
    chart_data.add_series("Amount", values)
    
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS, left, top, width, height, chart_data
    ).chart
    
    chart.has_legend = False
    chart.category_axis.tick_labels.font.size = Pt(8)
    chart.category_axis.tick_labels.font.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    chart.value_axis.tick_labels.font.size = Pt(7)
    chart.value_axis.tick_labels.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
    
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.font.size = Pt(7)
    plot.data_labels.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    
    series = plot.series[0]
    series.format.line.color.rgb = RGBColor(0x38, 0xBD, 0xF8)
    series.format.line.width = Pt(2.5)
    series.smooth = True
    
    return chart

def add_text_box(slide, text, left, top, width, height, size=10, bold=False, color=None, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_font(run, size, bold, color or RGBColor(0xFF, 0xFF, 0xFF))
    return txBox

def add_table(slide, headers, rows, left, top, width, row_height=Inches(0.28)):
    """Add a formatted table."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    col_w = width / n_cols
    
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, row_height * n_rows)
    table = table_shape.table
    
    # Headers
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = h
        set_font(run, 9, True, RGBColor(0x38, 0xBD, 0xF8))
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
    
    # Data rows
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(val)
            set_font(run, 8, False, RGBColor(0xE2, 0xE8, 0xF0))
    
    return table_shape

# ── Main ──

def generate(input_json):
    data = json.loads(input_json)
    role = data["role"]           # "invoice" | "payment"
    filter_label = data["filterLabel"]  # "Ekspor" | "Non-Ekspor"
    section = data["section"]
    
    role_label = "Invoice" if role == "invoice" else "Payment"
    date_label = today_label()
    
    prs = Presentation(TEMPLATE)
    delete_all_slides(prs)  # remove template's 5 sample slides; we build our own
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(
        prs,
        f"Laporan Monitoring {role_label}\n{filter_label}",
        f"per {date_label} | Periode: {section.get('latestPeriod', '-')}"
    )
    
    # Slide 2: Section break
    add_section_slide(prs, f"{role_label} {filter_label}\nKey Metrics")

    # Slide 2b: AI executive summary (only if narrative was generated)
    add_narrative_slide(prs, data.get("narrative"))

    # Slide 3: KPI cards
    slide = add_content_slide(prs, f"{role_label} {filter_label} — Ringkasan")
    
    total = section["totalAmount"]
    row_count = section["rowCount"]
    avg = section["averageAmount"]
    latest = section.get("latestPeriod", "-")
    
    # Determine overdue/high-risk
    status_mix = section.get("statusMix", [])
    if role == "invoice":
        overdue = sum(s["value"] for s in status_mix if s["label"].lower() != "current")
        risk_label = "Overdue (>Current)"
    else:
        overdue = sum(s["value"] for s in status_mix if "high risk" in s["label"].lower())
        risk_label = "High Risk Payment"
    
    # KPI boxes
    kpis = [
        ("Top Left", f"Total {role_label}", format_currency(total, True), f"{row_count:,} records"),
        ("Top Right", "Rata-rata", format_currency(avg, True), ""),
        ("Bottom Left", risk_label, format_currency(overdue, True), ""),
        ("Bottom Right", "Periode Terakhir", latest, ""),
    ]
    
    positions = [
        (Inches(0.5), Inches(1.4), Inches(3.0), Inches(1.3)),   # TL
        (Inches(3.7), Inches(1.4), Inches(3.0), Inches(1.3)),   # TR
        (Inches(7.0), Inches(1.4), Inches(3.0), Inches(1.3)),   # BL
        (Inches(10.2), Inches(1.4), Inches(2.8), Inches(1.3)),  # BR
    ]
    
    for (pos_name, label, value, sub), (x, y, w, h) in zip(kpis, positions):
        # Background rectangle
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE (rounded not available generically)
            x, y, w, h
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
        shape.line.fill.background()
        
        add_text_box(slide, label, x + Inches(0.15), y + Inches(0.1), w - Inches(0.3), Inches(0.25), 9, False, RGBColor(0x94, 0xA3, 0xB8))
        add_text_box(slide, value, x + Inches(0.15), y + Inches(0.4), w - Inches(0.3), Inches(0.4), 16, True, RGBColor(0xFF, 0xFF, 0xFF))
        if sub:
            add_text_box(slide, sub, x + Inches(0.15), y + Inches(0.85), w - Inches(0.3), Inches(0.2), 8, False, RGBColor(0x94, 0xA3, 0xB8))
    
    # Slide 4: Status/Aging Pie
    slide = add_content_slide(prs, f"{role_label} {filter_label} — Distribusi Status")
    
    status_order = (
        ["Current", "Bucket 1", "Bucket 2", "Bucket 3", "Bucket 4"]
        if role == "invoice"
        else ["No Risk", "Low Risk", "Warning", "Warning +", "High Risk", "High Risk +"]
    )
    
    status_data = {"labels": [], "values": []}
    for so in status_order:
        match = next((s for s in status_mix if s["label"].lower() == so.lower()), None)
        if match:
            status_data["labels"].append(match["label"])
            status_data["values"].append(match["value"])
    
    if status_data["labels"]:
        chart_title = "Status Tagihan (Aging Bucket)" if role == "invoice" else "Risk Status Distribution"
        add_text_box(slide, chart_title, Inches(0.3), Inches(1.2), Inches(5), Inches(0.25), 10, False, RGBColor(0x94, 0xA3, 0xB8))
        add_pie_chart(slide, chart_title, status_data["labels"], status_data["values"],
                      Inches(0.3), Inches(1.5), Inches(5.5), Inches(4.5))
        
        # Data table beside chart
        total_s = sum(status_data["values"]) or 1
        sorted_data = sorted(zip(status_data["labels"], status_data["values"]), key=lambda x: -x[1])
        
        headers = ["Kategori", "Nilai", "Share"]
        rows = [[l, format_currency(v, True), format_pct(v/total_s)] for l, v in sorted_data]
        add_table(slide, headers, rows, Inches(6.2), Inches(1.5), Inches(6.5))
    
    # Slide 5: Top Customers bar chart
    top_cust = section.get("topCustomers", [])[:8]
    if top_cust:
        slide = add_content_slide(prs, f"{role_label} {filter_label} — Top Customers")
        
        cust_labels = [c["label"] for c in top_cust]
        cust_values = [c["value"] for c in top_cust]
        
        add_text_box(slide, f"TOP {len(top_cust)} Customers", Inches(0.3), Inches(1.2), Inches(5), Inches(0.25), 10, False, RGBColor(0x94, 0xA3, 0xB8))
        add_bar_chart(slide, "", cust_labels, cust_values,
                      Inches(0.3), Inches(1.5), Inches(12.5), Inches(5.0))
    
    # Slide 6: Invoice Type pie (invoice only)
    if role == "invoice":
        inv_types = section.get("invoiceTypes", [])
        if inv_types:
            slide = add_content_slide(prs, f"{role_label} {filter_label} — Invoice Type")
            
            it_labels = [t["label"] for t in inv_types]
            it_values = [t["value"] for t in inv_types]
            
            add_text_box(slide, "Invoice Type Distribution", Inches(0.3), Inches(1.2), Inches(5), Inches(0.25), 10, False, RGBColor(0x94, 0xA3, 0xB8))
            add_pie_chart(slide, "", it_labels, it_values,
                          Inches(0.3), Inches(1.5), Inches(5.5), Inches(4.5))
            
            total_it = sum(it_values) or 1
            sorted_it = sorted(zip(it_labels, it_values), key=lambda x: -x[1])
            headers = ["Tipe Invoice", "Nilai", "Share"]
            rows = [[l, format_currency(v, True), format_pct(v/total_it)] for l, v in sorted_it]
            add_table(slide, headers, rows, Inches(6.2), Inches(1.5), Inches(6.5))
    
    # Slide 7: Monthly Trend line
    monthly = section.get("monthly", [])
    if len(monthly) > 1:
        slide = add_content_slide(prs, f"{role_label} {filter_label} — Tren Bulanan")
        
        m_labels = [m["label"] for m in monthly]
        m_values = [m["value"] for m in monthly]
        
        add_text_box(slide, "Tren Bulanan (Dalam Rupiah)", Inches(0.3), Inches(1.2), Inches(5), Inches(0.25), 10, False, RGBColor(0x94, 0xA3, 0xB8))
        add_line_chart(slide, "", m_labels, m_values,
                       Inches(0.3), Inches(1.5), Inches(12.5), Inches(5.0))
    
    # Slide 8: Top Records Table
    records = section.get("records", [])
    if records:
        slide = add_content_slide(prs, f"{role_label} {filter_label} — Detail Transaksi")
        
        top_records = sorted(records, key=lambda r: -r["amount"])[:15]
        headers = ["Customer", "Status", "Type", "Period", "Amount"]
        rows = [[
            r["customerName"],
            r["status"],
            r.get("invoiceType", r.get("customerType", "-")),
            r.get("periodLabel", "-"),
            format_currency(r["amount"]),
        ] for r in top_records]
        
        add_table(slide, headers, rows, Inches(0.3), Inches(1.3), Inches(12.8), Inches(0.25))
    
    # Slide 9: Kesimpulan
    if status_mix:
        slide = add_content_slide(prs, f"{role_label} {filter_label} — Kesimpulan")
        
        total_fmt = format_currency(total, True)
        top_status = status_mix[0] if status_mix else None
        
        conclusions = [
            f"Total {role_label} {filter_label}: {total_fmt} ({row_count:,} transaksi).",
            f"Periode terakhir: {latest}.",
        ]
        if top_status:
            conclusions.append(
                f"Status dominan: {top_status['label']} ({format_pct(top_status['share'])} — {format_currency(top_status['value'], True)})."
            )
        if overdue > 0:
            conclusions.append(f"{risk_label}: {format_currency(overdue, True)}.")
        conclusions.append(f"Rata-rata per transaksi: {format_currency(avg)}.")
        
        y = Inches(1.5)
        for line in conclusions:
            add_text_box(slide, line, Inches(0.5), y, Inches(12), Inches(0.35), 12, False, RGBColor(0xFF, 0xFF, 0xFF))
            y += Inches(0.4)
    
    # Slide 10: End
    add_end_slide(prs)
    
    return prs


if __name__ == "__main__":
    input_json = sys.stdin.read()
    prs = generate(input_json)
    output_path = sys.argv[1] if len(sys.argv) > 1 else "output.pptx"
    prs.save(output_path)
    print(f"OK:{output_path}")
